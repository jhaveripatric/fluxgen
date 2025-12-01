#!/usr/bin/env python3
"""
FluxGen EDO Pitch Deck Generator - FIXED VERSION
Generates professional PowerPoint presentation with speaker notes
FIXED: Proper speaker notes handling to avoid file corruption
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

# Helper function to safely add speaker notes
def add_speaker_notes(slide, notes_text):
    """Add speaker notes to a slide - with error handling"""
    try:
        notes_slide = slide.notes_slide
        text_frame = notes_slide.notes_text_frame
        text_frame.clear()  # Clear any existing notes
        text_frame.text = notes_text
    except Exception as e:
        print(f"   Warning: Could not add speaker notes - {str(e)}")
        # Continue anyway - slide will still be created

# FluxGen Brand Colors
DARK_BLUE = RGBColor(44, 62, 80)      # #2C3E50
COPPER = RGBColor(184, 115, 51)       # #B87333
SILVER = RGBColor(192, 192, 192)      # #C0C0C0
WHITE = RGBColor(255, 255, 255)

# Paths
BASE_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
LOGO_FULL = os.path.join(BASE_DIR, 'repo/app/static/images/fluxgen-logo-full.png')
LOGO_MONOGRAM = os.path.join(BASE_DIR, 'repo/app/static/images/fluxgen-logo-monogram.png')
OUTPUT_FILE = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'FluxGen-EDO-Presentation-COMPLETE.pptx')


def create_title_slide(prs):
    """Slide 1: Cover Slide"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BLUE
    
    # Logo at top center
    left = Inches(3.5)
    top = Inches(0.5)
    logo = slide.shapes.add_picture(LOGO_FULL, left, top, height=Inches(1.2))
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "FORGING TOMORROW'S WELDS"
    title_para = title_frame.paragraphs[0]
    title_para.alignment = PP_ALIGN.CENTER
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = COPPER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(2.8), Inches(8), Inches(1.2))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    
    subtitle_text = """FluxGen Industries Ltd.
Partnership Proposal to City of Airdrie Economic Development

Building Canada's First SAW Flux Manufacturing Capability"""
    
    subtitle_frame.text = subtitle_text
    for para in subtitle_frame.paragraphs:
        para.alignment = PP_ALIGN.CENTER
        para.font.size = Pt(18)
        para.font.color.rgb = SILVER
    
    # Separator line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(1.5), Inches(4.2), Inches(7), Inches(0.02)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = COPPER
    line.line.fill.background()
    
    # Presented to
    presented_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(0.8))
    presented_frame = presented_box.text_frame
    presented_text = """Presented to: Chad Sheldon, Economic Development Officer
City of Airdrie

Date: December 1, 2024"""
    presented_frame.text = presented_text
    for para in presented_frame.paragraphs:
        para.alignment = PP_ALIGN.CENTER
        para.font.size = Pt(14)
        para.font.color.rgb = SILVER
    
    # Team
    team_box = slide.shapes.add_textbox(Inches(1), Inches(5.8), Inches(8), Inches(1))
    team_frame = team_box.text_frame
    team_frame.text = "PRESENTING TEAM:"
    team_para = team_frame.paragraphs[0]
    team_para.alignment = PP_ALIGN.CENTER
    team_para.font.size = Pt(12)
    team_para.font.bold = True
    team_para.font.color.rgb = COPPER
    
    team_members = [
        "Pratik Jhaveri  |  Founder & Technical Director",
        "Arpan Patel     |  Chief Operating Officer",
        "Abhishek Patel  |  Chief Commercial Officer",
        "Bhargav Patel   |  Chief Compliance Officer"
    ]
    
    for member in team_members:
        p = team_frame.add_paragraph()
        p.text = member
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(11)
        p.font.color.rgb = SILVER
    
    # Footer
    footer_box = slide.shapes.add_textbox(Inches(1), Inches(7), Inches(8), Inches(0.3))
    footer_frame = footer_box.text_frame
    footer_frame.text = "FluxGen Industries Ltd. | Airdrie, Alberta"
    footer_para = footer_frame.paragraphs[0]
    footer_para.alignment = PP_ALIGN.CENTER
    footer_para.font.size = Pt(10)
    footer_para.font.color.rgb = SILVER
    
    # Speaker Notes
    notes = """OPENING (0:00 - 2:00 minutes)

"Chad, thanks so much for meeting with us today. I'm Pratik Jhaveri, Founder and Technical Director of FluxGen. With me are Arpan (COO), Abhishek (CCO), and Bhargav (Chief Compliance Officer). We're all equal partners in this venture.

We've prepared a short presentation — about 15 minutes — and then we'd love to hear your thoughts and answer questions. Sound good?

Perfect. So let me start with why we're here."""
    
    add_speaker_notes(slide, notes)
    return slide


def main():
    """Generate the complete EDO pitch deck with all 12 slides and speaker notes"""
    print("🔧 FluxGen EDO Pitch Deck Generator - FIXED VERSION")
    print("=" * 70)
    
    # Check logo files
    print(f"\n📁 Checking logo files...")
    if not os.path.exists(LOGO_FULL):
        print(f"❌ ERROR: Logo not found at {LOGO_FULL}")
        return
    if not os.path.exists(LOGO_MONOGRAM):
        print(f"❌ ERROR: Monogram logo not found at {LOGO_MONOGRAM}")
        return
    print("✅ Logo files found")
    
    # Create presentation
    print("\n📊 Creating presentation...")
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Test with just title slide first
    print("   Creating Slide 1: Cover Slide (TEST)...")
    try:
        create_title_slide(prs)
        print("   ✅ Slide 1 created successfully")
    except Exception as e:
        print(f"   ❌ ERROR creating slide: {str(e)}")
        return
    
    # Save presentation
    print(f"\n💾 Saving test presentation to: {OUTPUT_FILE}")
    try:
        prs.save(OUTPUT_FILE)
        print("✅ File saved successfully!")
    except Exception as e:
        print(f"❌ ERROR saving file: {str(e)}")
        return
    
    print("\n✅ SUCCESS!")
    print("=" * 70)
    print(f"📊 TEST presentation created: {OUTPUT_FILE}")
    print("\n🎯 NEXT STEP:")
    print("   Try opening this file. If it works, I'll generate the full 12-slide version.")
    print("\n" + "=" * 70)


if __name__ == "__main__":
    main()
