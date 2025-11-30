#!/usr/bin/env python3
"""
FluxGen EDO Pitch Deck Generator
Generates professional PowerPoint presentation for Economic Development Office meeting
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

# FluxGen Brand Colors
DARK_BLUE = RGBColor(44, 62, 80)      # #2C3E50
COPPER = RGBColor(184, 115, 51)       # #B87333
SILVER = RGBColor(192, 192, 192)      # #C0C0C0
WHITE = RGBColor(255, 255, 255)

# Paths
BASE_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
LOGO_FULL = os.path.join(BASE_DIR, 'repo/app/static/images/fluxgen-logo-full.png')
LOGO_MONOGRAM = os.path.join(BASE_DIR, 'repo/app/static/images/fluxgen-logo-monogram.png')
OUTPUT_FILE = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'FluxGen-EDO-Presentation.pptx')


def create_title_slide(prs):
    """Slide 1: Cover Slide"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BLUE
    
    # Logo at top center
    left = Inches(3.5)
    top = Inches(0.5)
    logo = slide.shapes.add_picture(LOGO_FULL, left, top, height=Inches(1.2))
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "FORGING TOMORROW'S WELDS"
    title_para = title_frame.paragraphs[0]
    title_para.alignment = PP_ALIGN.CENTER
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = COPPER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(2.8), Inches(8), Inches(1.2))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    
    subtitle_text = """FluxGen Industries Ltd.
Partnership Proposal to City of Airdrie Economic Development

Building Canada's First SAW Flux Manufacturing Capability"""
    
    subtitle_frame.text = subtitle_text
    for para in subtitle_frame.paragraphs:
        para.alignment = PP_ALIGN.CENTER
        para.font.size = Pt(18)
        para.font.color.rgb = SILVER
    
    # Separator line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(1.5), Inches(4.2), Inches(7), Inches(0.02)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = COPPER
    line.line.fill.background()
    
    # Presented to
    presented_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(0.8))
    presented_frame = presented_box.text_frame
    presented_text = """Presented to: Chad Sheldon, Economic Development Officer
City of Airdrie

Date: December 1, 2024"""
    presented_frame.text = presented_text
    for para in presented_frame.paragraphs:
        para.alignment = PP_ALIGN.CENTER
        para.font.size = Pt(14)
        para.font.color.rgb = SILVER
    
    # Team
    team_box = slide.shapes.add_textbox(Inches(1), Inches(5.8), Inches(8), Inches(1))
    team_frame = team_box.text_frame
    team_frame.text = "PRESENTING TEAM:"
    team_para = team_frame.paragraphs[0]
    team_para.alignment = PP_ALIGN.CENTER
    team_para.font.size = Pt(12)
    team_para.font.bold = True
    team_para.font.color.rgb = COPPER
    
    team_members = [
        "Pratik Jhaveri  |  Founder & Technical Director",
        "Arpan Patel     |  Chief Operating Officer",
        "Abhishek Patel  |  Chief Commercial Officer",
        "Bhargav Patel   |  Chief Compliance Officer"
    ]
    
    for member in team_members:
        p = team_frame.add_paragraph()
        p.text = member
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(11)
        p.font.color.rgb = SILVER
    
    # Footer
    footer_box = slide.shapes.add_textbox(Inches(1), Inches(7), Inches(8), Inches(0.3))
    footer_frame = footer_box.text_frame
    footer_frame.text = "FluxGen Industries Ltd. | Airdrie, Alberta"
    footer_para = footer_frame.paragraphs[0]
    footer_para.alignment = PP_ALIGN.CENTER
    footer_para.font.size = Pt(10)
    footer_para.font.color.rgb = SILVER


def create_opportunity_slide(prs):
    """Slide 2: The Opportunity"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BLUE
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
    title_frame = title_box.text_frame
    title_frame.text = "CANADA'S SAW FLUX MARKET: 100% IMPORT DEPENDENT"
    title_para = title_frame.paragraphs[0]
    title_para.alignment = PP_ALIGN.CENTER
    title_para.font.size = Pt(28)
    title_para.font.bold = True
    title_para.font.color.rgb = COPPER
    
    # Market box
    market_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(2), Inches(1.2), Inches(6), Inches(1.5)
    )
    market_box.fill.solid()
    market_box.fill.fore_color.rgb = RGBColor(52, 73, 94)  # Slightly lighter blue
    market_box.line.color.rgb = COPPER
    market_box.line.width = Pt(2)
    
    market_text = market_box.text_frame
    market_text.text = "$33 MILLION"
    market_para = market_text.paragraphs[0]
    market_para.alignment = PP_ALIGN.CENTER
    market_para.font.size = Pt(44)
    market_para.font.bold = True
    market_para.font.color.rgb = COPPER
    
    p2 = market_text.add_paragraph()
    p2.text = "Annual SAW Flux Imports"
    p2.alignment = PP_ALIGN.CENTER
    p2.font.size = Pt(18)
    p2.font.color.rgb = SILVER
    
    p3 = market_text.add_paragraph()
    p3.text = "11,000+ Tonnes Per Year"
    p3.alignment = PP_ALIGN.CENTER
    p3.font.size = Pt(16)
    p3.font.color.rgb = SILVER
    
    # Market Gap section
    gap_box = slide.shapes.add_textbox(Inches(0.8), Inches(3), Inches(8.4), Inches(1.8))
    gap_frame = gap_box.text_frame
    gap_frame.word_wrap = True
    
    gap_frame.text = "THE MARKET GAP:"
    gap_para = gap_frame.paragraphs[0]
    gap_para.font.size = Pt(18)
    gap_para.font.bold = True
    gap_para.font.color.rgb = COPPER
    gap_para.space_before = Pt(6)
    
    gap_points = [
        "• Zero domestic SAW flux manufacturers in Canada",
        "• 100% reliance on U.S., European, and Asian suppliers",
        "• Critical material for pipeline welding, shipbuilding, heavy equipment, and structural steel",
        "• Supply chain vulnerability — no Canadian backup source",
        "• Currency exchange risk and longer lead times"
    ]
    
    for point in gap_points:
        p = gap_frame.add_paragraph()
        p.text = point
        p.font.size = Pt(14)
        p.font.color.rgb = SILVER
        p.space_before = Pt(4)
    
    # Opportunity section
    opp_box = slide.shapes.add_textbox(Inches(0.8), Inches(5.2), Inches(8.4), Inches(1.2))
    opp_frame = opp_box.text_frame
    opp_frame.word_wrap = True
    
    opp_frame.text = "THE OPPORTUNITY:"
    opp_para = opp_frame.paragraphs[0]
    opp_para.font.size = Pt(18)
    opp_para.font.bold = True
    opp_para.font.color.rgb = COPPER
    
    p = opp_frame.add_paragraph()
    p.text = "Canada has the demand, the raw materials, and the skilled workforce."
    p.font.size = Pt(14)
    p.font.color.rgb = SILVER
    p.space_before = Pt(6)
    
    p2 = opp_frame.add_paragraph()
    p2.text = "What's missing? A domestic manufacturer."
    p2.font.size = Pt(14)
    p2.font.bold = True
    p2.font.color.rgb = SILVER
    
    p3 = opp_frame.add_paragraph()
    p3.text = "FluxGen is ready to fill that gap — and build it in Airdrie."
    p3.font.size = Pt(14)
    p3.font.bold = True
    p3.font.color.rgb = COPPER
    p3.space_before = Pt(6)
    
    # Logo footer
    slide.shapes.add_picture(LOGO_MONOGRAM, Inches(9.2), Inches(7), height=Inches(0.4))


def create_solution_slide(prs):
    """Slide 3: The Solution"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BLUE
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "FLUXGEN INDUSTRIES LTD."
    title_para = title_frame.paragraphs[0]
    title_para.alignment = PP_ALIGN.CENTER
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = COPPER
    
    subtitle = title_frame.add_paragraph()
    subtitle.text = "Canada's First SAW Flux & Alloy Manufacturer | Airdrie, Alberta"
    subtitle.alignment = PP_ALIGN.CENTER
    subtitle.font.size = Pt(16)
    subtitle.font.color.rgb = SILVER
    
    # Value Proposition boxes
    box_width = 2.6
    box_height = 1.8
    box_top = 1.3
    
    # Box 1: Proven Technology
    box1 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.7), Inches(box_top), Inches(box_width), Inches(box_height)
    )
    box1.fill.solid()
    box1.fill.fore_color.rgb = RGBColor(52, 73, 94)
    box1.line.color.rgb = COPPER
    box1.line.width = Pt(2)
    
    text1 = box1.text_frame
    text1.word_wrap = True
    text1.text = "PROVEN TECHNOLOGY"
    p1 = text1.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    p1.font.size = Pt(14)
    p1.font.bold = True
    p1.font.color.rgb = COPPER
    
    for point in ["Production-ready formulations", "25+ years family expertise", "NOT R&D — proven IP transfer"]:
        p = text1.add_paragraph()
        p.text = f"• {point}"
        p.font.size = Pt(11)
        p.font.color.rgb = SILVER
        p.space_before = Pt(4)
    
    # Box 2: Market Ready
    box2 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(3.7), Inches(box_top), Inches(box_width), Inches(box_height)
    )
    box2.fill.solid()
    box2.fill.fore_color.rgb = RGBColor(52, 73, 94)
    box2.line.color.rgb = COPPER
    box2.line.width = Pt(2)
    
    text2 = box2.text_frame
    text2.word_wrap = True
    text2.text = "MARKET READY"
    p2 = text2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    p2.font.size = Pt(14)
    p2.font.bold = True
    p2.font.color.rgb = COPPER
    
    for point in ["Immediate demand: 11,000 tonnes/year", "Zero domestic competition", "Pipeline welding = core Alberta market"]:
        p = text2.add_paragraph()
        p.text = f"• {point}"
        p.font.size = Pt(11)
        p.font.color.rgb = SILVER
        p.space_before = Pt(4)
    
    # Box 3: Export Potential
    box3 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(6.7), Inches(box_top), Inches(box_width), Inches(box_height)
    )
    box3.fill.solid()
    box3.fill.fore_color.rgb = RGBColor(52, 73, 94)
    box3.line.color.rgb = COPPER
    box3.line.width = Pt(2)
    
    text3 = box3.text_frame
    text3.word_wrap = True
    text3.text = "EXPORT POTENTIAL"
    p3 = text3.paragraphs[0]
    p3.alignment = PP_ALIGN.CENTER
    p3.font.size = Pt(14)
    p3.font.bold = True
    p3.font.color.rgb = COPPER
    
    for point in ["ArcNova Materials (export subsidiary)", "Target: U.S. & India markets", "Bring foreign revenue INTO Alberta"]:
        p = text3.add_paragraph()
        p.text = f"• {point}"
        p.font.size = Pt(11)
        p.font.color.rgb = SILVER
        p.space_before = Pt(4)
    
    # Phased Growth section
    growth_title = slide.shapes.add_textbox(Inches(0.7), Inches(3.4), Inches(8.6), Inches(0.4))
    growth_frame = growth_title.text_frame
    growth_frame.text = "PHASED GROWTH STRATEGY: De-Risked, Milestone-Based Expansion"
    growth_para = growth_frame.paragraphs[0]
    growth_para.alignment = PP_ALIGN.CENTER
    growth_para.font.size = Pt(16)
    growth_para.font.bold = True
    growth_para.font.color.rgb = COPPER
    
    # Phase boxes
    phase_width = 2.8
    phase_height = 2.2
    phase_top = 4
    
    phases = [
        {
            "title": "YEAR 1-2",
            "subtitle": "Pilot Production",
            "metrics": ["800 tonnes/year", "$2.5M revenue", "10 jobs created", "Canadian market focus"]
        },
        {
            "title": "YEAR 3-4",
            "subtitle": "Scale Production",
            "metrics": ["3,000 tonnes/year", "$12M revenue", "22 jobs created", "Expand market share"]
        },
        {
            "title": "YEAR 5-7",
            "subtitle": "Full Capacity + Exports",
            "metrics": ["6,000 tonnes/year", "$33M revenue", "40 jobs created", "Export + wire manufacturing"]
        }
    ]
    
    for i, phase in enumerate(phases):
        left = Inches(0.7 + (i * 3.1))
        phase_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left, Inches(phase_top), Inches(phase_width), Inches(phase_height)
        )
        phase_box.fill.solid()
        phase_box.fill.fore_color.rgb = RGBColor(62, 83, 104)
        phase_box.line.color.rgb = COPPER
        phase_box.line.width = Pt(2)
        
        text_frame = phase_box.text_frame
        text_frame.word_wrap = True
        text_frame.text = phase["title"]
        title_p = text_frame.paragraphs[0]
        title_p.alignment = PP_ALIGN.CENTER
        title_p.font.size = Pt(14)
        title_p.font.bold = True
        title_p.font.color.rgb = COPPER
        
        subtitle_p = text_frame.add_paragraph()
        subtitle_p.text = phase["subtitle"]
        subtitle_p.alignment = PP_ALIGN.CENTER
        subtitle_p.font.size = Pt(12)
        subtitle_p.font.bold = True
        subtitle_p.font.color.rgb = SILVER
        subtitle_p.space_before = Pt(2)
        
        for metric in phase["metrics"]:
            m = text_frame.add_paragraph()
            m.text = metric
            m.alignment = PP_ALIGN.CENTER
            m.font.size = Pt(10)
            m.font.color.rgb = SILVER
            m.space_before = Pt(4)
    
    # Logo footer
    slide.shapes.add_picture(LOGO_MONOGRAM, Inches(9.2), Inches(7), height=Inches(0.4))


def create_ask_slide(prs):
    """Slide 4: The Ask"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BLUE
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
    title_frame = title_box.text_frame
    title_frame.text = "PARTNERSHIP REQUEST: WHAT WE NEED FROM THE CITY OF AIRDRIE"
    title_para = title_frame.paragraphs[0]
    title_para.alignment = PP_ALIGN.CENTER
    title_para.font.size = Pt(22)
    title_para.font.bold = True
    title_para.font.color.rgb = COPPER
    
    # Content area
    content_y = 1.1
    
    asks = [
        {
            "title": "✓  INDUSTRIAL LAND: 10-15 ACRES",
            "details": [
                "Option A: Free land grant (community investment model)",
                "Option B: Heavily discounted long-term lease ($1/year for 10 years)"
            ]
        },
        {
            "title": "✓  LETTER OF MUNICIPAL SUPPORT",
            "details": [
                "For federal grant applications: PrairiesCan, Alberta Innovates, CanExport"
            ]
        },
        {
            "title": "✓  UTILITY CONNECTION ASSISTANCE",
            "details": [
                "Guidance on power, water, natural gas routing",
                "Introductions to utility providers (ENMAX, EPCOR, ATCO)"
            ]
        },
        {
            "title": "✓  PERMITTING COOPERATION",
            "details": [
                "Collaborative approach to building, environmental, and operational permits",
                "Streamlined approval process where possible"
            ]
        }
    ]
    
    for ask in asks:
        ask_box = slide.shapes.add_textbox(Inches(0.7), Inches(content_y), Inches(8.6), Inches(0.8))
        ask_frame = ask_box.text_frame
        ask_frame.word_wrap = True
        
        ask_frame.text = ask["title"]
        ask_para = ask_frame.paragraphs[0]
        ask_para.font.size = Pt(16)
        ask_para.font.bold = True
        ask_para.font.color.rgb = COPPER
        
        for detail in ask["details"]:
            p = ask_frame.add_paragraph()
            p.text = f"  {detail}"
            p.font.size = Pt(12)
            p.font.color.rgb = SILVER
            p.space_before = Pt(2)
        
        content_y += 0.95
    
    # What we're NOT asking for
    separator = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(1), Inches(5.2), Inches(8), Inches(0.02)
    )
    separator.fill.solid()
    separator.fill.fore_color.rgb = COPPER
    separator.line.fill.background()
    
    not_asking = slide.shapes.add_textbox(Inches(0.7), Inches(5.4), Inches(8.6), Inches(1.2))
    not_frame = not_asking.text_frame
    not_frame.word_wrap = True
    
    not_frame.text = "WHAT WE ARE NOT REQUESTING:"
    not_para = not_frame.paragraphs[0]
    not_para.font.size = Pt(14)
    not_para.font.bold = True
    not_para.font.color.rgb = COPPER
    
    not_items = [
        "✗  No cash grants or direct capital funding",
        "✗  No loan guarantees or financial backstops",
        "✗  No property tax abatements",
        "✗  No ongoing operational subsidies"
    ]
    
    for item in not_items:
        p = not_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(12)
        p.font.color.rgb = SILVER
        p.space_before = Pt(3)
    
    commitment = not_frame.add_paragraph()
    commitment.text = "\nWe will pay full property tax from Day 1. We will secure capital through federal/provincial grants and private investment."
    commitment.font.size = Pt(11)
    commitment.font.bold = True
    commitment.font.color.rgb = COPPER
    commitment.space_before = Pt(8)
    
    # Logo footer
    slide.shapes.add_picture(LOGO_MONOGRAM, Inches(9.2), Inches(7), height=Inches(0.4))


def create_roi_slide(prs):
    """Slide 5: Return on Investment for Airdrie"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BLUE
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
    title_frame = title_box.text_frame
    title_frame.text = "RETURN ON INVESTMENT FOR AIRDRIE"
    title_para = title_frame.paragraphs[0]
    title_para.alignment = PP_ALIGN.CENTER
    title_para.font.size = Pt(28)
    title_para.font.bold = True
    title_para.font.color.rgb = COPPER
    
    # Three-column layout for benefits
    col_width = 2.8
    col_height = 5
    col_top = 1.2
    
    columns = [
        {
            "title": "ECONOMIC IMPACT",
            "items": [
                "40 direct jobs by Year 7",
                "100+ indirect jobs (logistics, services, contractors)",
                "$500K-$1M annual property tax revenue",
                "$33M annual revenue at full capacity",
                "Attracts related industries (welding, fabrication)"
            ]
        },
        {
            "title": "STRATEGIC VALUE",
            "items": [
                "First-mover industrial leadership",
                "Tax base diversification beyond residential",
                "Advanced manufacturing sector anchor",
                "Export revenue generation",
                "Technology transfer and IP development"
            ]
        },
        {
            "title": "COMMUNITY BENEFITS",
            "items": [
                "Skilled trades employment",
                "Industrial land utilization",
                "Supply chain development",
                "Business services demand",
                "National visibility for Airdrie"
            ]
        }
    ]
    
    for i, col in enumerate(columns):
        left = Inches(0.7 + (i * 3.1))
        col_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left, Inches(col_top), Inches(col_width), Inches(col_height)
        )
        col_box.fill.solid()
        col_box.fill.fore_color.rgb = RGBColor(52, 73, 94)
        col_box.line.color.rgb = COPPER
        col_box.line.width = Pt(2)
        
        text_frame = col_box.text_frame
        text_frame.word_wrap = True
        text_frame.text = col["title"]
        title_p = text_frame.paragraphs[0]
        title_p.alignment = PP_ALIGN.CENTER
        title_p.font.size = Pt(14)
        title_p.font.bold = True
        title_p.font.color.rgb = COPPER
        
        for item in col["items"]:
            p = text_frame.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(11)
            p.font.color.rgb = SILVER
            p.space_before = Pt(6)
    
    # Bottom summary
    summary_box = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(8), Inches(0.8))
    summary_frame = summary_box.text_frame
    summary_frame.word_wrap = True
    summary_frame.text = "AIRDRIE'S INVESTMENT:"
    summary_para = summary_frame.paragraphs[0]
    summary_para.alignment = PP_ALIGN.CENTER
    summary_para.font.size = Pt(16)
    summary_para.font.bold = True
    summary_para.font.color.rgb = COPPER
    
    p2 = summary_frame.add_paragraph()
    p2.text = "Land commitment (retained ownership in lease model) + Municipal support letters"
    p2.alignment = PP_ALIGN.CENTER
    p2.font.size = Pt(13)
    p2.font.color.rgb = SILVER
    p2.space_before = Pt(4)
    
    p3 = summary_frame.add_paragraph()
    p3.text = "AIRDRIE'S RETURN: Generational industrial employer + diversified tax base + export economy participation"
    p3.alignment = PP_ALIGN.CENTER
    p3.font.size = Pt(13)
    p3.font.bold = True
    p3.font.color.rgb = COPPER
    p3.space_before = Pt(6)
    
    # Logo footer
    slide.shapes.add_picture(LOGO_MONOGRAM, Inches(9.2), Inches(7), height=Inches(0.4))


def create_team_slide(prs):
    """Slide 6: Our Team"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BLUE
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
    title_frame = title_box.text_frame
    title_frame.text = "OUR TEAM: PROVEN EXPERTISE"
    title_para = title_frame.paragraphs[0]
    title_para.alignment = PP_ALIGN.CENTER
    title_para.font.size = Pt(28)
    title_para.font.bold = True
    title_para.font.color.rgb = COPPER
    
    # Team member boxes (2x2 grid)
    team_members = [
        {
            "name": "Pratik Jhaveri",
            "title": "Founder & Technical Director",
            "highlights": [
                "Former owner of Jhaveri Weld Flux Ltd. (publicly listed Indian flux manufacturer)",
                "25+ years flux manufacturing expertise",
                "Complete formulation and process knowledge"
            ]
        },
        {
            "name": "Arpan Patel",
            "title": "Chief Operating Officer",
            "highlights": [
                "Logistics and supply chain background",
                "Loblaws Canada operations experience",
                "Manufacturing process optimization"
            ]
        },
        {
            "name": "Abhishek Patel",
            "title": "Chief Commercial Officer",
            "highlights": [
                "Sales and marketing expertise",
                "B2B industrial sales background",
                "Customer relationship management"
            ]
        },
        {
            "name": "Bhargav Patel",
            "title": "Chief Compliance Officer",
            "highlights": [
                "IT and systems background",
                "Detail-oriented compliance focus",
                "Regulatory and quality systems"
            ]
        }
    ]
    
    box_width = 4.2
    box_height = 2.6
    
    for i, member in enumerate(team_members):
        row = i // 2
        col = i % 2
        left = Inches(0.8 + (col * 4.6))
        top = Inches(1.2 + (row * 2.9))
        
        member_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left, top, Inches(box_width), Inches(box_height)
        )
        member_box.fill.solid()
        member_box.fill.fore_color.rgb = RGBColor(52, 73, 94)
        member_box.line.color.rgb = COPPER
        member_box.line.width = Pt(2)
        
        text_frame = member_box.text_frame
        text_frame.word_wrap = True
        text_frame.text = member["name"]
        name_p = text_frame.paragraphs[0]
        name_p.alignment = PP_ALIGN.CENTER
        name_p.font.size = Pt(16)
        name_p.font.bold = True
        name_p.font.color.rgb = COPPER
        
        title_p = text_frame.add_paragraph()
        title_p.text = member["title"]
        title_p.alignment = PP_ALIGN.CENTER
        title_p.font.size = Pt(13)
        title_p.font.color.rgb = SILVER
        title_p.space_before = Pt(2)
        
        for highlight in member["highlights"]:
            h = text_frame.add_paragraph()
            h.text = f"• {highlight}"
            h.font.size = Pt(10)
            h.font.color.rgb = SILVER
            h.space_before = Pt(4)
    
    # Logo footer
    slide.shapes.add_picture(LOGO_MONOGRAM, Inches(9.2), Inches(7), height=Inches(0.4))


def create_timeline_slide(prs):
    """Slide 7: Timeline & Milestones"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BLUE
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
    title_frame = title_box.text_frame
    title_frame.text = "TIMELINE & MILESTONES"
    title_para = title_frame.paragraphs[0]
    title_para.alignment = PP_ALIGN.CENTER
    title_para.font.size = Pt(28)
    title_para.font.bold = True
    title_para.font.color.rgb = COPPER
    
    # Timeline items
    timeline_items = [
        ("Q1 2025", "Land commitment + Federal grant applications"),
        ("Q2 2025", "Equipment procurement + Site development begins"),
        ("Q3-Q4 2025", "Building construction + Equipment installation"),
        ("Q1 2026", "AWS A5.17 certification + Staff hiring"),
        ("Q2 2027", "Production launch (800 tonnes Year 1)"),
        ("2028-2029", "Scale to 3,000 tonnes/year + Market expansion"),
        ("2030-2031", "Full capacity 6,000 tonnes + Export operations"),
        ("2032+", "Forward integration: Flux-cored wire manufacturing")
    ]
    
    y_position = 1.2
    for quarter, milestone in timeline_items:
        # Quarter box
        quarter_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.8), Inches(y_position), Inches(1.8), Inches(0.5)
        )
        quarter_box.fill.solid()
        quarter_box.fill.fore_color.rgb = COPPER
        quarter_box.line.fill.background()
        
        quarter_text = quarter_box.text_frame
        quarter_text.text = quarter
        quarter_p = quarter_text.paragraphs[0]
        quarter_p.alignment = PP_ALIGN.CENTER
        quarter_p.font.size = Pt(14)
        quarter_p.font.bold = True
        quarter_p.font.color.rgb = DARK_BLUE
        quarter_text.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        # Milestone text
        milestone_box = slide.shapes.add_textbox(Inches(2.8), Inches(y_position), Inches(6.4), Inches(0.5))
        milestone_frame = milestone_box.text_frame
        milestone_frame.text = milestone
        milestone_p = milestone_frame.paragraphs[0]
        milestone_p.font.size = Pt(13)
        milestone_p.font.color.rgb = SILVER
        milestone_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        y_position += 0.65
    
    # Logo footer
    slide.shapes.add_picture(LOGO_MONOGRAM, Inches(9.2), Inches(7), height=Inches(0.4))


def create_why_airdrie_slide(prs):
    """Slide 8: Why Airdrie?"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BLUE
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
    title_frame = title_box.text_frame
    title_frame.text = "WHY AIRDRIE?"
    title_para = title_frame.paragraphs[0]
    title_para.alignment = PP_ALIGN.CENTER
    title_para.font.size = Pt(28)
    title_para.font.bold = True
    title_para.font.color.rgb = COPPER
    
    # Reasons
    reasons = [
        {
            "title": "STRATEGIC LOCATION",
            "points": [
                "Calgary proximity — access to skilled workforce and suppliers",
                "Transportation infrastructure — Highway 2, CP Rail access",
                "Natural gas availability for kilns and dryers"
            ]
        },
        {
            "title": "ECONOMIC DEVELOPMENT ALIGNMENT",
            "points": [
                "Airdrie 2024 Economic Strategy targets advanced manufacturing",
                "Tax base diversification beyond residential growth",
                "Proven track record of supporting industrial development"
            ]
        },
        {
            "title": "INDUSTRIAL LAND AVAILABILITY",
            "points": [
                "East Lake Industrial area has suitable parcels",
                "Zoning already in place for manufacturing",
                "Room for phased expansion through 2031"
            ]
        },
        {
            "title": "BUSINESS-FRIENDLY ENVIRONMENT",
            "points": [
                "Responsive EDO team (Chad's proactive approach)",
                "Municipal willingness to partner with growth companies",
                "Alberta's competitive business tax environment"
            ]
        }
    ]
    
    y_pos = 1.2
    for reason in reasons:
        reason_box = slide.shapes.add_textbox(Inches(0.8), Inches(y_pos), Inches(8.4), Inches(1.2))
        reason_frame = reason_box.text_frame
        reason_frame.word_wrap = True
        
        reason_frame.text = reason["title"]
        title_p = reason_frame.paragraphs[0]
        title_p.font.size = Pt(16)
        title_p.font.bold = True
        title_p.font.color.rgb = COPPER
        
        for point in reason["points"]:
            p = reason_frame.add_paragraph()
            p.text = f"• {point}"
            p.font.size = Pt(12)
            p.font.color.rgb = SILVER
            p.space_before = Pt(3)
        
        y_pos += 1.35
    
    # Logo footer
    slide.shapes.add_picture(LOGO_MONOGRAM, Inches(9.2), Inches(7), height=Inches(0.4))


def create_next_steps_slide(prs):
    """Slide 9: Next Steps"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BLUE
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
    title_frame = title_box.text_frame
    title_frame.text = "NEXT STEPS"
    title_para = title_frame.paragraphs[0]
    title_para.alignment = PP_ALIGN.CENTER
    title_para.font.size = Pt(28)
    title_para.font.bold = True
    title_para.font.color.rgb = COPPER
    
    # Steps with timeline
    steps = [
        ("TODAY", "Initial conversation and partnership discussion", COPPER),
        ("WEEK 1", "City identifies potential industrial parcels (10-15 acres)", SILVER),
        ("WEEK 2", "FluxGen conducts site feasibility assessments", SILVER),
        ("WEEK 3-4", "Draft preliminary MOU (land terms, timeline, commitments)", SILVER),
        ("JANUARY 2025", "City provides municipal support letter for federal grants", COPPER),
        ("JANUARY 2025", "FluxGen submits PrairiesCan and Alberta Innovates applications", SILVER),
        ("Q2 2025", "Federal grant decisions + Equipment procurement begins", SILVER),
        ("Q3 2025", "Site development and building construction kickoff", COPPER)
    ]
    
    y_pos = 1.2
    for when, what, color in steps:
        # Timeline marker
        marker = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.8), Inches(y_pos), Inches(1.5), Inches(0.4)
        )
        marker.fill.solid()
        marker.fill.fore_color.rgb = color
        marker.line.fill.background()
        
        marker_text = marker.text_frame
        marker_text.text = when
        marker_p = marker_text.paragraphs[0]
        marker_p.alignment = PP_ALIGN.CENTER
        marker_p.font.size = Pt(12)
        marker_p.font.bold = True
        marker_p.font.color.rgb = DARK_BLUE
        marker_text.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        # Step description
        step_box = slide.shapes.add_textbox(Inches(2.5), Inches(y_pos), Inches(6.7), Inches(0.4))
        step_frame = step_box.text_frame
        step_frame.text = what
        step_p = step_frame.paragraphs[0]
        step_p.font.size = Pt(12)
        step_p.font.color.rgb = SILVER
        step_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        y_pos += 0.55
    
    # Bottom call to action
    cta_box = slide.shapes.add_textbox(Inches(1), Inches(6.2), Inches(8), Inches(0.8))
    cta_frame = cta_box.text_frame
    cta_frame.word_wrap = True
    cta_frame.text = "OUR ASK TODAY:"
    cta_para = cta_frame.paragraphs[0]
    cta_para.alignment = PP_ALIGN.CENTER
    cta_para.font.size = Pt(18)
    cta_para.font.bold = True
    cta_para.font.color.rgb = COPPER
    
    p2 = cta_frame.add_paragraph()
    p2.text = "Commitment to explore land partnership + Municipal support for federal grant applications"
    p2.alignment = PP_ALIGN.CENTER
    p2.font.size = Pt(14)
    p2.font.color.rgb = SILVER
    p2.space_before = Pt(6)
    
    # Logo footer
    slide.shapes.add_picture(LOGO_MONOGRAM, Inches(9.2), Inches(7), height=Inches(0.4))


def create_closing_slide(prs):
    """Slide 10: Closing / Thank You"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BLUE
    
    # Logo at center
    logo = slide.shapes.add_picture(LOGO_FULL, Inches(3), Inches(1.5), height=Inches(1.5))
    
    # Main message
    message_box = slide.shapes.add_textbox(Inches(1), Inches(3.2), Inches(8), Inches(1.5))
    message_frame = message_box.text_frame
    message_frame.word_wrap = True
    message_frame.text = "FORGING TOMORROW'S WELDS"
    msg_para = message_frame.paragraphs[0]
    msg_para.alignment = PP_ALIGN.CENTER
    msg_para.font.size = Pt(36)
    msg_para.font.bold = True
    msg_para.font.color.rgb = COPPER
    
    p2 = message_frame.add_paragraph()
    p2.text = "Let's Build Canada's First SAW Flux Manufacturing Capability"
    p2.alignment = PP_ALIGN.CENTER
    p2.font.size = Pt(20)
    p2.font.color.rgb = SILVER
    p2.space_before = Pt(12)
    
    p3 = message_frame.add_paragraph()
    p3.text = "Together in Airdrie, Alberta"
    p3.alignment = PP_ALIGN.CENTER
    p3.font.size = Pt(18)
    p3.font.color.rgb = SILVER
    p3.space_before = Pt(8)
    
    # Contact info
    contact_box = slide.shapes.add_textbox(Inches(1), Inches(5.2), Inches(8), Inches(1.2))
    contact_frame = contact_box.text_frame
    contact_frame.word_wrap = True
    
    contact_frame.text = "CONTACT:"
    contact_para = contact_frame.paragraphs[0]
    contact_para.alignment = PP_ALIGN.CENTER
    contact_para.font.size = Pt(14)
    contact_para.font.bold = True
    contact_para.font.color.rgb = COPPER
    
    contacts = [
        "Pratik Jhaveri  |  Founder & Technical Director",
        "pratik@fluxgen.ca  |  (403) XXX-XXXX",
        "",
        "FluxGen Industries Ltd.",
        "Airdrie, Alberta, Canada"
    ]
    
    for contact in contacts:
        p = contact_frame.add_paragraph()
        p.text = contact
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(12)
        p.font.color.rgb = SILVER
        p.space_before = Pt(3)
    
    # Thank you
    thanks_box = slide.shapes.add_textbox(Inches(1), Inches(6.8), Inches(8), Inches(0.4))
    thanks_frame = thanks_box.text_frame
    thanks_frame.text = "Thank you for your time and consideration"
    thanks_para = thanks_frame.paragraphs[0]
    thanks_para.alignment = PP_ALIGN.CENTER
    thanks_para.font.size = Pt(16)
    thanks_para.font.bold = True
    thanks_para.font.color.rgb = COPPER


def main():
    """Generate the complete EDO pitch deck"""
    print("🔧 FluxGen EDO Pitch Deck Generator")
    print("=" * 60)
    
    # Check logo files
    print(f"\n📁 Checking logo files...")
    if not os.path.exists(LOGO_FULL):
        print(f"❌ ERROR: Logo not found at {LOGO_FULL}")
        return
    if not os.path.exists(LOGO_MONOGRAM):
        print(f"❌ ERROR: Monogram logo not found at {LOGO_MONOGRAM}")
        return
    print("✅ Logo files found")
    
    # Create presentation
    print("\n📊 Creating presentation...")
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Generate all slides
    slides_created = []
    
    print("   Creating Slide 1: Cover Slide...")
    create_title_slide(prs)
    slides_created.append("Cover Slide")
    
    print("   Creating Slide 2: The Opportunity...")
    create_opportunity_slide(prs)
    slides_created.append("The Opportunity")
    
    print("   Creating Slide 3: The Solution...")
    create_solution_slide(prs)
    slides_created.append("The Solution")
    
    print("   Creating Slide 4: The Ask...")
    create_ask_slide(prs)
    slides_created.append("The Ask")
    
    print("   Creating Slide 5: ROI for Airdrie...")
    create_roi_slide(prs)
    slides_created.append("Return on Investment")
    
    print("   Creating Slide 6: Our Team...")
    create_team_slide(prs)
    slides_created.append("Our Team")
    
    print("   Creating Slide 7: Timeline & Milestones...")
    create_timeline_slide(prs)
    slides_created.append("Timeline & Milestones")
    
    print("   Creating Slide 8: Why Airdrie?...")
    create_why_airdrie_slide(prs)
    slides_created.append("Why Airdrie?")
    
    print("   Creating Slide 9: Next Steps...")
    create_next_steps_slide(prs)
    slides_created.append("Next Steps")
    
    print("   Creating Slide 10: Closing/Thank You...")
    create_closing_slide(prs)
    slides_created.append("Closing")
    
    # Save presentation
    print(f"\n💾 Saving presentation to: {OUTPUT_FILE}")
    prs.save(OUTPUT_FILE)
    
    print("\n✅ SUCCESS!")
    print("=" * 60)
    print(f"📊 Presentation created: {OUTPUT_FILE}")
    print(f"📝 Total slides: {len(slides_created)}")
    print("\n📋 Slides included:")
    for i, slide in enumerate(slides_created, 1):
        print(f"   {i}. {slide}")
    
    print("\n🎯 NEXT STEPS:")
    print("   1. Open the .pptx file to review")
    print("   2. Upload to Google Slides (File → Import slides)")
    print("   3. Make any final adjustments")
    print("   4. Present with confidence!")
    print("\n" + "=" * 60)


if __name__ == "__main__":
    main()
